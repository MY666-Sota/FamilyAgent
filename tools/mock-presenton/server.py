"""
mock-presenton — 轻量 Presenton 上游模拟服务（端口 7860，可改）
================================================================
用途：真镜像的 LLM 生成依赖云端 API（受 SSL 断流阻塞），无法在受限网络下跑通
      端到端。此 mock 完整模拟真实 Presenton 的 API 契约（端点/请求体/响应/
      文件下载），用 python-pptx 本地生成真实 .pptx，让 presenton-mcp 的层B
      e2e 在不依赖云端 LLM 的情况下验证整条链路：
          presenton-mcp(9002) → 本服务(7860) → 生成 .pptx → MCP 下载落 shared/outputs/

模拟的真实契约（2026-06-19 实测 ghcr.nju.edu.cn/presenton/presenton:latest）：
  POST /api/v1/ppt/presentation/generate
    请求体: {content, n_slides?, language?, template?, export_as?, instructions?}
    响应:   {presentation_id, path: "/app_data/.../xxx.pptx", edit_path}
  GET  /app_data/{...}.pptx   → 静态下载生成的文件
  GET  /                       → 健康检查（真实服务根路径返回 200）

启动：
  pip install python-pptx fastapi uvicorn
  python tools/mock-presenton/server.py            # 默认 7860
  PORT=7861 python tools/mock-presenton/server.py  # 自定义端口

注意：仅用于联调验证，生成的是占位 PPT（标题+要点文本），无 AI 内容、无配图。
      真镜像网络通畅后应切回真服务（presenton-mcp 无需改动，API 契约一致）。
"""
import os
import re
import uuid
from pathlib import Path

import uvicorn
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel

# mock 生成的文件落在本服务自己的 app_data 目录（模拟容器的 /app_data 卷）
APP_DATA_DIR = Path(__file__).parent / "_app_data"
APP_DATA_DIR.mkdir(parents=True, exist_ok=True)

app = FastAPI(title="Mock Presenton", version="1.0.0")


class GenerateRequest(BaseModel):
    content: str
    slides_markdown: list[str] | None = None
    instructions: str | None = None
    n_slides: int | None = None
    language: str | None = None
    template: str = "general"
    export_as: str = "pptx"


def _build_pptx(content: str, n_slides: int, dest: Path) -> int:
    """用 python-pptx 生成占位 PPT。返回实际页数。"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    # 标题页：取 content 第一行作主标题
    lines = [l.strip() for l in content.splitlines() if l.strip()]
    main_title = lines[0] if lines else "Presentation"
    title_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = main_title
    if len(s0.placeholders) > 1:
        s0.placeholders[1].text = "由 mock-presenton 生成（联调占位）"

    # 内容页：把 content 切成 n_slides 段
    body_layout = prs.slide_layouts[1]
    # 优先按 markdown ## 切分
    md_sections = re.split(r"\n(?=##\s)", content)
    if len(md_sections) > 1:
        chunks = md_sections
    else:
        # 否则平均分配剩余行
        rest = lines[1:] if len(lines) > 1 else lines
        per = max(1, len(rest) // max(1, n_slides))
        chunks = [rest[i:i + per] for i in range(0, len(rest), per)] or [["内容"]]
        chunks = ["\n".join(c) if isinstance(c, list) else c for c in chunks]

    made = 1
    for i, chunk in enumerate(chunks[:n_slides] if n_slides else chunks):
        s = prs.slides.add_slide(body_layout)
        chunk_lines = [l.strip("# ").strip() for l in str(chunk).splitlines() if l.strip()]
        s.shapes.title.text = chunk_lines[0] if chunk_lines else f"第 {i+1} 页"
        body = s.placeholders[1].text_frame
        body.text = ""
        for bl in (chunk_lines[1:] or ["（要点占位）"]):
            p = body.add_paragraph()
            p.text = bl
            p.font.size = Pt(18)
        made += 1

    prs.save(str(dest))
    return made


@app.get("/")
def root():
    """真实服务根路径返回 200（Next.js 前端）。"""
    return {"status": "ok", "service": "mock-presenton"}


@app.post("/api/v1/ppt/presentation/generate")
def generate(req: GenerateRequest):
    """模拟同步生成端点，返回真实契约的 path/edit_path。"""
    if req.export_as != "pptx":
        raise HTTPException(status_code=400, detail="mock 仅支持 export_as=pptx")

    pid = uuid.uuid4()
    # 模拟真实目录结构：/app_data/{presentation_id}/xxx.pptx
    sub = APP_DATA_DIR / str(pid)
    sub.mkdir(parents=True, exist_ok=True)
    out = sub / "presentation.pptx"

    n = req.n_slides or 3
    made = _build_pptx(req.content, n, out)

    rel_path = f"/app_data/{pid}/presentation.pptx"
    return {
        "presentation_id": str(pid),
        "path": rel_path,
        "edit_path": f"/presentation?id={pid}",
        "n_slides": made,
    }


@app.get("/app_data/{pid}/{name}")
def download(pid: str, name: str):
    """静态下载生成的 pptx（模拟容器的 /app_data StaticFiles 挂载）。"""
    # 防目录穿越
    target = (APP_DATA_DIR / pid / name).resolve()
    if not str(target).startswith(str(APP_DATA_DIR.resolve())):
        raise HTTPException(status_code=400, detail="非法路径")
    if not target.exists():
        raise HTTPException(status_code=404, detail="文件不存在")
    return FileResponse(str(target))


if __name__ == "__main__":
    port = int(os.getenv("PORT", "7860"))
    uvicorn.run(app, host="0.0.0.0", port=port)